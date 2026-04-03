#!/usr/bin/env python3
"""Generate portfolio presentation PowerPoint from Markdown content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Color palette ---
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # deep navy
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)    # brand blue
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)    # teal accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUBTLE_BG = RGBColor(0xF5, 0xF7, 0xFA)      # light slide bg

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bold_prefix=True, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        # Handle bold prefix (text before " -- ")
        if bold_prefix and " -- " in item:
            prefix, rest = item.split(" -- ", 1)
            run = p.add_run()
            run.text = prefix + " -- "
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = color
            run.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = color
            run.font.name = "Calibri"

    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None):
    """Add a formatted table. rows_data[0] is header row."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.4 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_TEXT

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Accent bar at top
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_TEAL)

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
            "Sales Engineering Portfolio", font_size=44, bold=True,
            color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.6),
            "Dwayne Dreakford  |  Sales Engineer", font_size=24,
            color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
            "linkedin.com/in/dwaynedreakford   |   github.com/ddreakford",
            font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# ============================================================
# SLIDE 2: Introductions & Agenda (dark scheme)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_TEAL)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Introductions & Agenda", font_size=32, bold=True, color=WHITE)

# Left column: About Me
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "About Me", font_size=22, bold=True, color=ACCENT_TEAL)

# About Me - main points (bold)
about_me_main = [
    "Sales Engineer",
    "My primary focus is on selling",
    "My \"1a\" focus is on force multiplication through sharing of best practices and reusable assets that make SEs more effective",
    "I approach every technical investment through a lens of sales impact:",
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.5), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(about_me_main):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = item
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = "Calibri"

# About Me - sub-questions (not bold, indented feel)
about_me_sub = [
    "How will this facilitate selling conversations?",
    "Shorten deal cycles? Improve win rates?",
    "Scale beyond me?",
]
add_bullet_list(slide, Inches(1.2), Inches(4.2), Inches(5.1), Inches(1.5),
                about_me_sub, font_size=15, color=LIGHT_GRAY, bold_prefix=False, spacing=Pt(4))

# Right column: Agenda
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
            "Agenda", font_size=22, bold=True, color=ACCENT_TEAL)

agenda_items = [
    "1.  SE Philosophy",
    "2.  Portfolio Overview",
    "3.  Containerized Appium Testing Environment",
    "4.  Mobile Application Security Assessment Framework",
    "5.  Feature Project: Multi Test Type Demo Sandbox",
    "6.  Demo",
    "7.  Cross-Project Themes & Summary",
]
add_bullet_list(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.5),
                agenda_items, font_size=16, color=LIGHT_GRAY, bold_prefix=False, spacing=Pt(8))

add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# ============================================================
# SLIDE 3: SE Philosophy (dark scheme)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_TEAL)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "SE Philosophy", font_size=32, bold=True, color=WHITE)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "How I Think About Sales Engineering", font_size=22, bold=True, color=ACCENT_TEAL)

add_textbox(slide, Inches(0.8), Inches(2.1), Inches(11), Inches(0.5),
            "Every project in this portfolio was built to answer four questions:",
            font_size=16, color=LIGHT_GRAY)

philosophy_items = [
    "Understand the domain and challenges -- to maximizing business value",
    "Identify key requirements and practical value -- of addressing the challenges",
    "Clarify how you, in differentiated fashion, provide a solution -- that meets the requirements",
    "Prove the viability of the business case -- as it pertains to the solution to be provided",
]
add_bullet_list(slide, Inches(1.2), Inches(2.7), Inches(10.5), Inches(2.5),
                philosophy_items, font_size=18, color=LIGHT_GRAY, spacing=Pt(12))

add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
            "The Common Thread", font_size=22, bold=True, color=ACCENT_TEAL)

add_bullet_list(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(1.5), [
    "Few things inspire confidence like tangible examples.",
    "These projects provide tangible examples -- reusable assets that compound in value across engagements, reduce marginal effort over time, and enable the SE team to focus on customer conversations rather than setup and configuration.",
], font_size=16, color=LIGHT_GRAY, bold_prefix=False, spacing=Pt(8))

add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# ============================================================
# SLIDE 4: Portfolio Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Portfolio Overview", font_size=32, bold=True, color=DARK_BG)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "Three Projects, Similar Intent: Facilitate Clarification of Value",
            font_size=20, bold=True, color=ACCENT_BLUE)

add_table(slide, Inches(0.8), Inches(1.9), Inches(11.5), [
    ["Project", "Domain", "Core SE Value"],
    ["Multi Test Type Demo Sandbox", "Automated and manual tests",
     "Realistic AUT; API + UI + CI + reporting in one project"],
    ["Containerized Appium Testing Environment", "Mobile test automation",
     "Reduced demo/POC setup from days to hours; multi-persona engagement"],
    ["Mobile Application Security Assessment Framework", "AI-assisted / hybrid and manual workflows",
     "Packaged know-how into a repeatable process; unblocked opportunities"],
], col_widths=[Inches(3.8), Inches(3.2), Inches(4.5)])

add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
            "Consistent Design Principles", font_size=20, bold=True, color=ACCENT_BLUE)

add_bullet_list(slide, Inches(1.2), Inches(4.8), Inches(10.5), Inches(2.5), [
    "Ready-to-run -- command line ready; sensible defaults, configuration available",
    "Flexible, Reproducible deployment -- eliminate environment variability",
    "Documented for reuse -- tutorials, guides and templates so others can adopt",
    "Audience suitable outputs -- reports and dashboards for stakeholders, code & config for engineers",
], font_size=15, spacing=Pt(8))

# ============================================================
# SLIDE 5: Containerized Appium Testing Environment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Containerized Appium Testing Environment", font_size=32, bold=True, color=DARK_BG)

# Challenge
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(3.5), Inches(0.5),
            "Challenge", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(3.5), Inches(2.0), [
    "SEs lost days to environment setup before demos",
    "Demos constrained or skipped due to config complexity",
    "Evaluations dragged on",
], font_size=14, bold_prefix=False, spacing=Pt(6))

# Solution
add_textbox(slide, Inches(4.8), Inches(1.3), Inches(4.2), Inches(0.5),
            "Solution", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(4.8), Inches(1.9), Inches(4.2), Inches(3.0), [
    "CLI-triggered scenarios -- basic functionality, custom scenarios, advanced capabilities (self-healing, accessibility, performance testing)",
    "Docker Build and Compose -- eliminate dependency issues, plug in to CI without mods",
    "Published as community code -- public consumption, feedback and market awareness",
], font_size=14, spacing=Pt(6))

# Impact
add_textbox(slide, Inches(9.5), Inches(1.3), Inches(3.3), Inches(0.5),
            "Impact", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(9.5), Inches(1.9), Inches(3.3), Inches(3.0), [
    "Demo and POC setup reduced from days to hours",
    "Multi-persona engagement in a single session",
    "SEs deliver advanced demos within hours of receiving cloud access",
    "Competitive differentiation through live demonstration of advanced capabilities",
], font_size=14, bold_prefix=False, spacing=Pt(6))

# Tech bar
add_shape_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5), RGBColor(0xF0, 0xF4, 0xF8))
add_textbox(slide, Inches(1.0), Inches(5.85), Inches(11), Inches(0.4),
            "Tech: Docker, Gradle, TestNG, Appium, Java, Python",
            font_size=13, color=MED_GRAY)

# ============================================================
# SLIDE 6: Mobile Application Security Assessment Framework
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Mobile Application Security Assessment Framework", font_size=32, bold=True, color=DARK_BG)

# Challenge
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(3.5), Inches(0.5),
            "Challenge", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(3.5), Inches(2.0), [
    "Pipeline progression gated by 1-2 specialists",
    "Deals stalled waiting for assessments",
    "Ad hoc process; 8-12 hours per assessment",
], font_size=14, bold_prefix=False, spacing=Pt(6))

# Solution
add_textbox(slide, Inches(4.8), Inches(1.3), Inches(4.2), Inches(0.5),
            "Solution", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(4.8), Inches(1.9), Inches(4.2), Inches(3.0), [
    "Phase-by-phase workflow -- setup, recon, vulnerability hunting, MASVS-RESILIENCE analysis, reporting",
    "Three approaches -- AI-assisted (fastest), manual (most thorough), hybrid (recommended)",
    "CVSS v3.1 scoring -- with professional report templates",
    "Comparative analysis -- track remediation effectiveness across app versions",
    "Multi-audience deliverables -- technical reports, executive summaries, stakeholder email templates",
], font_size=14, spacing=Pt(6))

# Impact
add_textbox(slide, Inches(9.5), Inches(1.3), Inches(3.3), Inches(0.5),
            "Impact", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(9.5), Inches(1.9), Inches(3.3), Inches(3.0), [
    "Assessment time reduced from 8-12 hours to 4-6 hours",
    "Security conversations no longer bottlenecked by specialist availability",
    "Discovery conversations uncovered requirements competitors missed",
    "Comparative analysis encouraged ongoing engagement",
], font_size=14, bold_prefix=False, spacing=Pt(6))

# Tech bar
add_shape_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5), RGBColor(0xF0, 0xF4, 0xF8))
add_textbox(slide, Inches(1.0), Inches(5.85), Inches(11), Inches(0.4),
            "Tech: JADX, APKTool, Python, Bash, OWASP MASVS v2.0, CVSS v3.1, Claude Code",
            font_size=13, color=MED_GRAY)

# ============================================================
# SLIDE 7: Feature Project - Multi Test Type Demo Sandbox
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.4),
            "FEATURE PROJECT", font_size=14, bold=True, color=ACCENT_TEAL)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
            "Multi Test Type Demo Sandbox", font_size=32, bold=True, color=DARK_BG)

# The Problem
add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5),
            "The Problem", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.0), [
    "Most demos showcase capabilities in isolation",
    "No end-to-end reference -- how do API tests, UI tests, reporting and CI fit together?",
    "RCA is under-represented -- what happens when tests fail?",
    "What about manual testers? -- how and when should manual tests be automated?",
], font_size=15, spacing=Pt(6))

# Architecture table
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
            "The Solution", font_size=20, bold=True, color=ACCENT_BLUE)

add_table(slide, Inches(0.8), Inches(4.8), Inches(11.5), [
    ["Layer", "Implementation"],
    ["System Under Test", "Restful-Booker Platform (multi-service Docker deployed app)"],
    ["API Testing", "RestAssured + TestNG (auth, booking CRUD)"],
    ["UI Testing", "Selenium WebDriver + TestNG (admin login, homepage, booking flow, contact form)"],
    ["Reporting", "Allure Report (dashboards, step-level detail, screenshots, trends)"],
    ["Manual Testing", "Step-by-step tutorial with reference screenshots"],
    ["CI/CD", "GitHub Actions (full suite + RCA demo + Allure artifact upload)"],
], col_widths=[Inches(2.5), Inches(9.0)])

# ============================================================
# SLIDE 8: Execution, Reporting & RCA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Multi Test Type Demo: Execution, Reporting & RCA",
            font_size=28, bold=True, color=DARK_BG)

# Test Coverage table
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
            "Test Coverage", font_size=20, bold=True, color=ACCENT_BLUE)

add_table(slide, Inches(0.8), Inches(1.8), Inches(6.0), [
    ["Suite", "Tests", "What It Covers"],
    ["API Tests (Auth, Booking CRUD)", "5", "Cookie-based auth, CRUD via REST"],
    ["UI Test (Admin Login)", "1", "Browser-based admin panel login"],
    ["Regression Suite", "17", "Full UI validation (maps 1:1 to manual)"],
    ["RCA Demo", "2", "Intentional failures for RCA walkthrough"],
], col_widths=[Inches(2.5), Inches(0.7), Inches(2.8)])

# RCA section
add_textbox(slide, Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.5),
            "RCA Demo: Teaching Through Failure", font_size=18, bold=True, color=ACCENT_BLUE)

add_bullet_list(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(2.5), [
    "Two tests designed to fail -- both are test defects, not system defects",
    "RCA-001 (API): Asserts HTTP 200 but service correctly returns 201 Created",
    "RCA-002 (UI): Asserts wrong page title",
    "Allure dashboard walkthrough shows how to identify, classify and document failures",
], font_size=14, spacing=Pt(6))

# Manual + Automated
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
            "Manual + Automated, Side by Side", font_size=18, bold=True, color=ACCENT_BLUE)

add_bullet_list(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5), [
    "Three manual test scenarios map 1:1 to the automated regression suite",
    "Speaks to teams in transition and concretely demonstrates the automation value proposition",
], font_size=15, bold_prefix=False, spacing=Pt(6))

# ============================================================
# SLIDE 9: Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_TEAL)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Demo: Multi Test Type Demo Sandbox", font_size=32, bold=True, color=WHITE)

demo_steps = [
    "1.  Clone and start -- git clone --recurse-submodules + docker compose up",
    "2.  Run the automated suite -- ./gradlew clean test (23 tests)",
    "3.  View Allure Report -- dashboards, test detail, screenshots",
    "4.  Run the RCA demo -- ./gradlew rcaDemo (intentional failures)",
    "5.  RCA walkthrough -- investigate failures in Allure",
    "6.  Manual test tutorial -- reference screenshots alongside automated equivalents",
]
add_bullet_list(slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(4.0),
                demo_steps, font_size=18, color=LIGHT_GRAY, spacing=Pt(14))

add_textbox(slide, Inches(1.2), Inches(6.0), Inches(10), Inches(0.5),
            "[Live demo, recorded walkthrough, or screenshot-based -- adapt to venue]",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# ============================================================
# SLIDE 10: Cross-Project Themes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Cross-Project Themes", font_size=32, bold=True, color=DARK_BG)

# Column 1
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(3.6), Inches(0.5),
            "Engineering Excellence in Sales", font_size=16, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(3.6), Inches(3.5), [
    "Build reusable assets that compound -- each framework serves multiple engagements and reduces marginal effort",
    "Apply software engineering best practices to SE deliverables -- version control, CI/CD, containerization, documentation",
    "Balance technical depth with business applicability -- every project includes architecture docs, getting-started guides and outcome narratives alongside the code",
], font_size=13, spacing=Pt(6))

# Column 2
add_textbox(slide, Inches(5.0), Inches(1.4), Inches(3.6), Inches(0.5),
            "Strategic Technical Leadership", font_size=16, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.5), [
    "Identify gaps in the SE toolkit -- and proactively build frameworks to address them",
    "Design for extensibility -- modular architecture so others can build on the foundation",
    "Contribute to community knowledge -- open-source where possible to strengthen the ecosystem",
], font_size=13, spacing=Pt(6))

# Column 3
add_textbox(slide, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.5),
            "Business Impact Orientation", font_size=16, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(9.2), Inches(1.9), Inches(3.6), Inches(3.5), [
    "Every technical investment is evaluated through a sales impact lens -- Will this shorten deal cycles? Enable new conversations? Improve win rates?",
    "Quantify SE impact in business terms -- setup time reduction, assessment throughput, multi-persona engagement",
    "Build bridges from technical capabilities to business outcomes",
], font_size=13, spacing=Pt(6))

# ============================================================
# SLIDE 11: Summary & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_TEAL)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Summary & Next Steps", font_size=32, bold=True, color=WHITE)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(0.5),
            "Key Takeaways", font_size=20, bold=True, color=ACCENT_TEAL)

takeaways = [
    "Technical assets should be built with sales impact in mind -- not just \"what does it do\" but \"how does it help us sell\"",
    "Ready-to-run beats configurable -- SEs need to focus on customer conversations, not framework setup",
    "Codified methodology scales; individual expertise doesn't -- frameworks turn specialist skills into team capabilities",
    "Show, don't tell -- live demonstrations, working code and professional reports build confidence that slides cannot",
]
add_bullet_list(slide, Inches(1.0), Inches(2.1), Inches(10.5), Inches(3.0),
                takeaways, font_size=16, color=LIGHT_GRAY, spacing=Pt(10))

add_textbox(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(0.5),
            "Explore the Portfolio", font_size=20, bold=True, color=ACCENT_TEAL)

links = [
    "Portfolio: github.com/ddreakford/sales-engineering-portfolio",
    "Multi Test Type Demo: github.com/ddreakford/multi-test-type-demo",
    "Containerized Appium: github.com/ddreakford/CommunityCode-AppiumCodeExamples",
]
add_bullet_list(slide, Inches(1.0), Inches(5.3), Inches(7), Inches(1.5),
                links, font_size=14, color=MED_GRAY, bold_prefix=False, spacing=Pt(4))

add_textbox(slide, Inches(8.5), Inches(4.8), Inches(4), Inches(0.5),
            "Get in Touch", font_size=20, bold=True, color=ACCENT_TEAL)

contact = [
    "LinkedIn: linkedin.com/in/dwaynedreakford",
    "GitHub: github.com/ddreakford",
]
add_bullet_list(slide, Inches(8.7), Inches(5.3), Inches(4), Inches(1.0),
                contact, font_size=14, color=MED_GRAY, bold_prefix=False, spacing=Pt(4))

add_textbox(slide, Inches(1.5), Inches(6.6), Inches(10), Inches(0.5),
            "Thank you -- questions?", font_size=24, color=WHITE,
            alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# ============================================================
# Save
# ============================================================
output_path = "presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
